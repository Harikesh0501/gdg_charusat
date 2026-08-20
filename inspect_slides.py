import pptx

prs = pptx.Presentation('SIH2026-IDEA-Presentation-Format (1).pptx')
print(f"Total Slides: {len(prs.slides)}")

for idx, slide in enumerate(prs.slides):
    print(f"\n==================== SLIDE {idx + 1} ====================")
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            print(f"  Shape [{shape.name}]:\n    {repr(txt)}\n")
